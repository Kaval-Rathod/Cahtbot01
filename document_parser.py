import os
from pptx import Presentation
from docx import Document
import PyPDF2

def read_text_file(filepath):
    with open(filepath, 'r') as file:
        return file.read()

def read_ppt_file(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_docx_file(filepath):
    doc = Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def read_pdf_file(filepath):
    reader = PyPDF2.PdfReader(filepath)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def parse_document(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.txt':
        return read_text_file(filepath)
    elif ext == '.pptx':
        return read_ppt_file(filepath)
    elif ext == '.docx':
        return read_docx_file(filepath)
    elif ext == '.pdf':
        return read_pdf_file(filepath)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
